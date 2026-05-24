"""CV text parser with string matching algorithms."""

import re
from dataclasses import dataclass, field

# Skill taxonomy
SKILL_DB: list[str] = [
    # Frontend
    "React", "Vue.js", "Angular", "Next.js", "TypeScript", "JavaScript",
    "HTML", "CSS", "Tailwind CSS", "Bootstrap", "SASS", "Astro",
    # Backend
    "Node.js", "Express", "Hono", "NestJS", "Fastify",
    "Python", "Django", "Flask", "FastAPI",
    "Go", "Fiber", "Gin",
    "Java", "Spring Boot", "Kotlin",
    "PHP", "Laravel", "CodeIgniter",
    "Ruby", "Ruby on Rails",
    "Rust", "C++", "C#", ".NET",
    # Mobile
    "React Native", "Flutter", "Swift", "Dart", "Expo",
    # Database
    "PostgreSQL", "MySQL", "MongoDB", "Redis", "SQLite", "SQL",
    # DevOps / Cloud
    "Docker", "Kubernetes", "AWS", "GCP", "Azure",
    "CI/CD", "GitHub Actions", "Terraform", "Linux",
    # MLOps / ML Infrastructure
    "MLflow", "Kubeflow", "KServe", "Feast", "MLOps",
    # Design
    "Figma", "Adobe XD", "Sketch", "Photoshop", "Illustrator",
    # AI / ML / Data Science Frameworks
    "TensorFlow", "PyTorch", "Pandas", "Scikit-learn",
    "XGBoost", "CatBoost", "LightGBM", "Random Forest",
    "Transformers", "Hugging Face", "LangChain", "FAISS",
    "Streamlit", "Gradio",
    # Neural Network Architectures
    "ANN", "CNN", "LSTM", "RNN", "FFNN",
    # ML Algorithms
    "KNN", "GNB", "Simulated Annealing", "Genetic Algorithm",
    # Data & Analytics
    "R", "Tableau", "Spreadsheet", "NumPy", "Matplotlib", "Seaborn",
    "Generative AI", "LLM",
    # APIs & Protocols
    "REST API", "GraphQL", "gRPC",
    # Version Control & Tools
    "Git", "Jira", "Agile", "Scrum",
]

# Skill aliases for fuzzy matching
SKILL_ALIASES: dict[str, str] = {
    "reactjs": "React", "react.js": "React",
    "vuejs": "Vue.js", "vue": "Vue.js",
    "angularjs": "Angular",
    "ts": "TypeScript", "typescript": "TypeScript",
    "js": "JavaScript", "javascript": "JavaScript",
    "nodejs": "Node.js", "node": "Node.js",
    "expressjs": "Express", "express.js": "Express",
    "golang": "Go",
    "postgres": "PostgreSQL", "psql": "PostgreSQL",
    "mongo": "MongoDB",
    "k8s": "Kubernetes",
    "tailwind": "Tailwind CSS",
    "nextjs": "Next.js", "next": "Next.js",
    "nestjs": "NestJS",
    "springboot": "Spring Boot",
    "ruby on rails": "Ruby on Rails", "rails": "Ruby on Rails",
    "react native": "React Native",
    "github actions": "GitHub Actions",
    "rest": "REST API", "restful": "REST API", "restapi": "REST API",
    "graphql": "GraphQL",
    "tensorflow": "TensorFlow",
    "pytorch": "PyTorch",
    "scikit-learn": "Scikit-learn", "sklearn": "Scikit-learn", "scikit learn": "Scikit-learn",
    "ci/cd": "CI/CD", "cicd": "CI/CD",
    "xgboost": "XGBoost",
    "catboost": "CatBoost",
    "lightgbm": "LightGBM",
    "random forest": "Random Forest", "randomforest": "Random Forest",
    "langchain": "LangChain", "lang chain": "LangChain",
    "faiss": "FAISS",
    "streamlit": "Streamlit",
    "gradio": "Gradio",
    "transformers": "Transformers",
    "hugging face": "Hugging Face", "huggingface": "Hugging Face",
    "generative ai": "Generative AI", "generativeai": "Generative AI",
    "large language model": "LLM",
    "mlflow": "MLflow",
    "kubeflow": "Kubeflow",
    "kserve": "KServe",
    "feast": "Feast",
    "mlops": "MLOps",
    "simulated annealing": "Simulated Annealing",
    "genetic algorithm": "Genetic Algorithm",
    "knn": "KNN", "k-nearest neighbor": "KNN",
    "gnb": "GNB", "gaussian naive bayes": "GNB",
    "astro": "Astro",
    "tableau": "Tableau",
}


def levenshtein_distance(s1: str, s2: str) -> int:
    """Levenshtein distance for fuzzy matching."""
    if len(s1) < len(s2):
        return levenshtein_distance(s2, s1)
    if len(s2) == 0:
        return len(s1)
    prev = list(range(len(s2) + 1))
    for i, c1 in enumerate(s1):
        curr = [i + 1]
        for j, c2 in enumerate(s2):
            cost = 0 if c1 == c2 else 1
            curr.append(min(curr[j] + 1, prev[j + 1] + 1, prev[j] + cost))
        prev = curr
    return prev[-1]


class AhoCorasick:
    """Multi-pattern matching."""

    def __init__(self) -> None:
        self.goto: list[dict[str, int]] = [{}]
        self.fail: list[int] = [0]
        self.output: list[list[str]] = [[]]

    def add_pattern(self, pattern: str, label: str) -> None:
        state = 0
        for ch in pattern.lower():
            if ch not in self.goto[state]:
                self.goto[state][ch] = len(self.goto)
                self.goto.append({})
                self.fail.append(0)
                self.output.append([])
            state = self.goto[state][ch]
        self.output[state].append(label)

    def build(self) -> None:
        from collections import deque
        q: deque[int] = deque()
        for ch, s in self.goto[0].items():
            q.append(s)
        while q:
            r = q.popleft()
            for ch, s in self.goto[r].items():
                q.append(s)
                state = self.fail[r]
                while state != 0 and ch not in self.goto[state]:
                    state = self.fail[state]
                self.fail[s] = self.goto[state].get(ch, 0)
                if self.fail[s] == s:
                    self.fail[s] = 0
                self.output[s] = self.output[s] + self.output[self.fail[s]]

    def search(self, text: str) -> set[str]:
        results: set[str] = set()
        state = 0
        for ch in text.lower():
            while state != 0 and ch not in self.goto[state]:
                state = self.fail[state]
            state = self.goto[state].get(ch, 0)
            for label in self.output[state]:
                results.add(label)
        return results


def build_skill_matcher() -> AhoCorasick:
    """Build Aho-Corasick automaton for skills."""
    ac = AhoCorasick()
    for skill in SKILL_DB:
        ac.add_pattern(skill.lower(), skill)
    for alias, canonical in SKILL_ALIASES.items():
        ac.add_pattern(alias, canonical)
    ac.build()
    return ac


_SKILL_MATCHER = build_skill_matcher()


def extract_skills_from_text(text: str) -> list[str]:
    """Extract skills using Aho-Corasick + Levenshtein fallback."""
    # Phase 1: Aho-Corasick exact + alias matching
    found = _SKILL_MATCHER.search(text)

    # Phase 2: Levenshtein fuzzy match on remaining words
    words = re.findall(r'\b[A-Za-z][A-Za-z.#+/\-]{1,20}\b', text)
    for word in words:
        w_lower = word.lower()
        if any(w_lower == s.lower() for s in found):
            continue
        for skill in SKILL_DB:
            if levenshtein_distance(w_lower, skill.lower()) <= 2 and len(skill) > 3:
                found.add(skill)
                break

    return sorted(found)


def extract_text(file_bytes: bytes, file_type: str) -> str:
    """Extract raw text from a document given its bytes and type."""
    import tempfile
    from pathlib import Path

    ext = file_type.lower()
    text = ""

    try:
        if ext == "pdf":
            try:
                import pypdfium2 as pdfium

                pdf = pdfium.PdfDocument(file_bytes)
                pages = []
                for page in pdf:
                    textpage = page.get_textpage()
                    pages.append(textpage.get_text_bounded())
                    textpage.close()
                    page.close()
                pdf.close()
                text = "\n".join(pages)
            except Exception:
                text = file_bytes.decode("utf-8", errors="ignore")
        elif ext in ("docx", "doc"):
            try:
                import docx

                with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
                    tmp.write(file_bytes)
                    tmp_path = tmp.name
                doc = docx.Document(tmp_path)
                text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                Path(tmp_path).unlink(missing_ok=True)
            except Exception:
                text = file_bytes.decode("utf-8", errors="ignore")
        elif ext == "pptx":
            try:
                from pptx import Presentation

                with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                    tmp.write(file_bytes)
                    tmp_path = tmp.name
                prs = Presentation(tmp_path)
                parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            parts.append(shape.text)
                text = "\n".join(parts)
                Path(tmp_path).unlink(missing_ok=True)
            except Exception:
                text = file_bytes.decode("utf-8", errors="ignore")
        else:
            text = file_bytes.decode("utf-8", errors="ignore")
    except Exception:
        text = file_bytes.decode("utf-8", errors="ignore")

    return text


@dataclass
class ParsedCV:
    """Structured CV data."""
    name: str = ""
    email: str = ""
    phone: str = ""
    summary: str = ""
    skills: list[str] = field(default_factory=list)
    education: list[dict] = field(default_factory=list)
    experience: list[dict] = field(default_factory=list)
    organizational_experience: list[dict] = field(default_factory=list)
    projects: list[dict] = field(default_factory=list)
    certifications: list[dict] = field(default_factory=list)
    portfolio_urls: list[str] = field(default_factory=list)
    years_of_experience: int | None = None


_SECTION_HEADERS = re.compile(
    r'^(?:'
    r'educations?|work\s+experiences?|organizational\s+experiences?|'
    r'achievements?|volunteering|projects?|certifications?|'
    r'skills?|awards?|publications?|languages?|interests?|references?|'
    r'pengalaman\s+(?:kerja|organisasi)?|pendidikan|proyek|sertifikasi|keahlian'
    r')$',
    re.IGNORECASE,
)


def _split_sections(text: str) -> dict[str, str]:
    """Split CV text into named sections by header detection."""
    sections: dict[str, str] = {}
    current_key = "header"
    current_lines: list[str] = []

    for line in text.split('\n'):
        stripped = line.strip()
        if stripped and _SECTION_HEADERS.match(stripped):
            sections[current_key] = '\n'.join(current_lines).strip()
            current_key = stripped.lower()
            current_lines = []
        else:
            current_lines.append(line)

    sections[current_key] = '\n'.join(current_lines).strip()
    return sections


def _normalize_section_key(key: str) -> str:
    """Map section key variants to canonical names."""
    k = key.lower().strip()
    if re.match(r'education', k):
        return 'education'
    if re.match(r'work\s*exp', k):
        return 'work_experience'
    if re.match(r'organizational\s*exp', k) or re.match(r'pengalaman\s+organisasi', k):
        return 'org_experience'
    if re.match(r'project', k) or re.match(r'proyek', k):
        return 'projects'
    if re.match(r'certif', k) or re.match(r'sertif', k):
        return 'certifications'
    if re.match(r'achieve', k) or re.match(r'award', k):
        return 'achievements'
    if re.match(r'voluntee', k):
        return 'volunteering'
    return k


def extract_emails(text: str) -> list[str]:
    return re.findall(r'[\w.+-]+@[\w-]+\.[\w.-]+', text)


def extract_phones(text: str) -> list[str]:
    return re.findall(r'(?:\+62|62|0)\d[\d\s\-]{8,14}', text)


def extract_urls(text: str) -> list[str]:
    urls = re.findall(r'https?://[^\s<>"\']+', text)
    portfolio_domains = ['github.com', 'linkedin.com', 'dribbble.com', 'behance.net', 'gitlab.com']
    return [u for u in urls if any(d in u for d in portfolio_domains)]


def extract_name_heuristic(text: str) -> str:
    """Extract name from first non-empty line that looks like a person's name."""
    for line in text.strip().split('\n'):
        line = line.strip()
        if 2 < len(line) < 60 and not re.search(r'[@\d|]', line) and not line.startswith('http'):
            # Reject lines that look like section headers or institutions
            if not _SECTION_HEADERS.match(line):
                return line
    return ""


def _parse_education_section(text: str) -> list[dict]:
    """Parse education section into structured entries."""
    entries = []
    # Split by blank lines or date patterns indicating new entries
    chunks = re.split(r'\n{2,}', text.strip())
    for chunk in chunks:
        if not chunk.strip():
            continue
        lines = [l.strip() for l in chunk.split('\n') if l.strip()]
        if not lines:
            continue
        entry: dict = {}
        # First line often is institution name
        entry['university'] = lines[0]
        for line in lines:
            # GPA
            gpa_m = re.search(r'GPA[:\s]+(\d+\.?\d*)', line, re.IGNORECASE)
            if gpa_m:
                entry['gpa'] = gpa_m.group(1)
            # Degree
            deg_m = re.search(
                r'(Bachelor|Master|PhD|S1|S2|S3|Diploma)[^\n]*?(of|in|:)?\s*([A-Z][^\n,]+)',
                line, re.IGNORECASE,
            )
            if deg_m:
                entry.setdefault('major', deg_m.group(0).strip()[:80])
            # Dates
            date_m = re.search(r'(\w+\s+\d{4})\s*[–\-]\s*(\w+\s+\d{4}|\w+)', line)
            if date_m:
                entry.setdefault('start', date_m.group(1))
                entry.setdefault('end', date_m.group(2))
        entries.append(entry)
    return entries


def _parse_experience_entries(text: str) -> list[dict]:
    """Parse work/org experience sections into structured entries."""
    entries = []
    # Split by company/org patterns (lines with ' – ' or dates)
    chunks = re.split(r'\n{2,}', text.strip())
    for chunk in chunks:
        lines = [l.strip() for l in chunk.split('\n') if l.strip()]
        if not lines or len(lines) < 1:
            continue
        entry: dict = {}
        header_line = lines[0]
        # Pattern: "Company – Role (type) | Location"
        header_m = re.match(r'^(.+?)\s*[–\-]\s*(.+?)(?:\s*\|\s*(.+))?$', header_line)
        if header_m:
            entry['company'] = header_m.group(1).strip()
            role_part = header_m.group(2).strip()
            # Strip parens like "(Contract, Remote)"
            entry['position'] = re.sub(r'\s*\([^)]*\)', '', role_part).strip()
        else:
            entry['company'] = header_line[:80]

        # Date line often second line
        for line in lines[1:3]:
            date_m = re.search(
                r'(\w+\s+\d{4})\s*[–\-]\s*(\w+\s+\d{4}|\w+\s+\d{4}|\w+)',
                line,
            )
            if date_m:
                entry['start'] = date_m.group(1)
                entry['end'] = date_m.group(2)
                break

        # Bullet points as description
        bullets = [l.lstrip('•·-').strip() for l in lines if l.startswith(('•', '·', '-', '*'))]
        if bullets:
            entry['description'] = ' '.join(bullets)[:500]

        if entry:
            entries.append(entry)
    return entries


def _parse_projects_section(text: str) -> list[dict]:
    """Parse projects section.

    Format observed in CVs:
      Title | Tech1, Tech2, Tech3 | URL1 | URL2   Date
      • bullet description
    """
    projects = []
    current: dict | None = None

    for line in text.split('\n'):
        stripped = line.strip()
        if not stripped:
            continue

        # Project header: contains ' | ' separating title from tech stack
        if '|' in stripped and not stripped.startswith(('•', '·', '-', '*')):
            if current:
                projects.append(current)
            parts = [p.strip() for p in stripped.split('|')]
            title_date = parts[0]
            # Title may have trailing date — strip it
            title = re.sub(r'\s+(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)\w*\s+\d{4}$', '', title_date).strip()
            tech_part = parts[1] if len(parts) > 1 else ""
            tech_stack = [t.strip() for t in re.split(r'[,&]+', tech_part) if t.strip()]
            urls = [p for p in parts[2:] if p.lower().startswith(('http', 'github', 'repository', 'repo'))]
            current = {"title": title, "tech_stack": tech_stack, "url": urls[0] if urls else "", "description": ""}
        elif stripped.startswith(('•', '·', '-', '*')) and current is not None:
            bullet = stripped.lstrip('•·-* ').strip()
            current['description'] = (current.get('description', '') + ' ' + bullet).strip()[:400]
        else:
            # Could be a date line or continuation — skip
            pass

    if current:
        projects.append(current)
    return projects


def _parse_certifications_section(text: str) -> list[dict]:
    """Parse certifications section.

    Format: • Cert Name | Tech1, Tech2 | URL   or   Cert Name | Tech1 | URL
    Also extract tech skills from cert tags.
    """
    certs = []
    for line in text.split('\n'):
        stripped = line.strip().lstrip('•·-* ')
        if not stripped:
            continue
        if '|' in stripped:
            parts = [p.strip() for p in stripped.split('|')]
            name = parts[0]
            tech_part = parts[1] if len(parts) > 1 else ""
            # Issuer often embedded in name: "IBM AI Engineering" → issuer=IBM
            issuer_m = re.match(r'^(IBM|Google|AWS|Microsoft|Meta|Oracle|Coursera|Udemy|LinkedIn)', name, re.IGNORECASE)
            certs.append({
                "name": name,
                "issuer": issuer_m.group(1) if issuer_m else "",
                "tech_tags": [t.strip() for t in re.split(r'[,&]+', tech_part) if t.strip()],
            })
        elif stripped:
            issuer_m = re.match(r'^(IBM|Google|AWS|Microsoft|Meta|Oracle|Coursera|Udemy|LinkedIn)', stripped, re.IGNORECASE)
            certs.append({
                "name": stripped[:120],
                "issuer": issuer_m.group(1) if issuer_m else "",
                "tech_tags": [],
            })
    return certs


def _extract_summary(header_text: str) -> str:
    """Extract summary paragraph from CV header block."""
    lines = header_text.split('\n')
    summary_lines = []
    in_summary = False
    for line in lines:
        stripped = line.strip()
        # Skip header metadata (email, phone, URLs, name)
        if not stripped:
            if in_summary and summary_lines:
                break  # end of summary paragraph
            continue
        if re.search(r'[@|]|https?://', stripped):
            continue
        if _SECTION_HEADERS.match(stripped):
            break
        # First substantial paragraph after metadata is the summary
        if len(stripped) > 40:
            in_summary = True
            summary_lines.append(stripped)
    return ' '.join(summary_lines)[:600]


def parse_cv_text(text: str) -> ParsedCV:
    """Section-aware CV parser. Primary extraction path for fallback when LLM unavailable."""
    result = ParsedCV()

    # Global extractions
    result.name = extract_name_heuristic(text)
    emails = extract_emails(text)
    if emails:
        result.email = emails[0]
    phones = extract_phones(text)
    if phones:
        result.phone = phones[0].strip()
    result.portfolio_urls = extract_urls(text)

    # Split into sections
    raw_sections = _split_sections(text)
    canonical: dict[str, str] = {}
    for k, v in raw_sections.items():
        canonical[_normalize_section_key(k)] = v

    # Header section → summary
    header_text = canonical.get('header', raw_sections.get('header', ''))
    result.summary = _extract_summary(header_text)

    # Education
    edu_text = canonical.get('education', '')
    if edu_text:
        result.education = _parse_education_section(edu_text)

    # Work experience
    work_text = canonical.get('work_experience', '')
    if work_text:
        result.experience = _parse_experience_entries(work_text)

    # Work experience year count
    if result.experience:
        years: set[int] = set()
        for exp in result.experience:
            for yr_str in [exp.get('start', ''), exp.get('end', '')]:
                m = re.search(r'\d{4}', yr_str or '')
                if m:
                    years.add(int(m.group()))
        if len(years) >= 2:
            result.years_of_experience = max(years) - min(years)

    # Organizational experience
    org_text = canonical.get('org_experience', '')
    if org_text:
        result.organizational_experience = _parse_experience_entries(org_text)

    # Include achievements/volunteering in org_experience
    for extra_key in ('achievements', 'volunteering'):
        extra_text = canonical.get(extra_key, '')
        if extra_text:
            result.organizational_experience.extend(_parse_experience_entries(extra_text))

    # Projects
    proj_text = canonical.get('projects', '')
    if proj_text:
        result.projects = _parse_projects_section(proj_text)

    # Certifications
    cert_text = canonical.get('certifications', '')
    if cert_text:
        result.certifications = _parse_certifications_section(cert_text)

    # Skills: scan ALL sections with Aho-Corasick
    # Also extract tech tags from projects and certifications
    extra_skill_tokens: list[str] = []
    for proj in result.projects:
        extra_skill_tokens.extend(proj.get('tech_stack', []))
    for cert in result.certifications:
        extra_skill_tokens.extend(cert.get('tech_tags', []))
    full_scan_text = text + '\n' + ' '.join(extra_skill_tokens)
    result.skills = extract_skills_from_text(full_scan_text)

    return result
